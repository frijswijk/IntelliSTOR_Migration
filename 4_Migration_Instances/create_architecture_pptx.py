#!/usr/bin/env python3
"""Create IntelliSTOR Architecture PowerPoint with visual diagrams."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MED_BLUE = RGBColor(0x2E, 0x6B, 0x9E)
LIGHT_BLUE = RGBColor(0x5B, 0xA0, 0xD0)
ACCENT_GREEN = RGBColor(0x2D, 0x8E, 0x4E)
ACCENT_ORANGE = RGBColor(0xD4, 0x7B, 0x2A)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_PURPLE = RGBColor(0x7B, 0x41, 0xA0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
MED_GRAY = RGBColor(0x80, 0x80, 0x80)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Title background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.5), Inches(3), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_ORANGE
    bar.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "IntelliSTOR Architecture"
    p.font.size = Pt(44)
    p.font.color.rgb = WHITE
    p.font.bold = True

    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = "Sections, Indexing & Extraction"
    p2.font.size = Pt(28)
    p2.font.color.rgb = LIGHT_BLUE

    # Description
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11), Inches(2))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p3 = tf2.paragraphs[0]
    p3.text = "How Section Security, MAP File Indexing, and Extraction Patterns\nwork together in the IntelliSTOR document management system"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(0xB0, 0xC4, 0xDE)


def add_box(slide, left, top, width, height, fill_color, text, font_size=11, font_color=WHITE, bold=False, align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(0x20, 0x20, 0x20)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = align
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = align
    return shape


def add_arrow(slide, start_left, start_top, end_left, end_top):
    """Add a connector line (simple line shape)."""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        Emu(start_left), Emu(start_top),
        Emu(end_left), Emu(end_top)
    )
    connector.line.color.rgb = MED_GRAY
    connector.line.width = Pt(2)
    return connector


def add_text(slide, left, top, width, height, text, font_size=11, font_color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = align
    return txBox


def add_overview_slide(prs):
    """Slide 2: Three systems overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "Three Independent Systems — One Integrated Workflow", 28, DARK_BLUE, True)

    # Accent line
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(12), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_ORANGE
    bar.line.fill.background()

    # Three system boxes
    y = Inches(1.4)
    box_h = Inches(1.8)
    box_w = Inches(3.7)

    # Box 1: Section Security
    add_box(slide, Inches(0.5), y, box_w, box_h, MED_BLUE,
            "SECTION SECURITY\n(STYPE)", 16, WHITE, True)
    add_text(slide, Inches(0.7), y + Inches(1.0), box_w - Inches(0.4), Inches(0.8),
             "Who can see which pages?\n• STYPE_SECTION ACL per section\n• RPT SECTIONHDR → page ranges\n• Per-user page filtering", 10, WHITE)

    # Box 2: MAP Indexing
    add_box(slide, Inches(4.7), y, box_w, box_h, ACCENT_GREEN,
            "MAP FILE INDEXING", 16, WHITE, True)
    add_text(slide, Inches(4.9), y + Inches(1.0), box_w - Inches(0.4), Inches(0.8),
             "Find pages by field value\n• IS_INDEXED=1 fields in MAP\n• Binary search → page numbers\n• Report-wide scope", 10, WHITE)

    # Box 3: Extraction
    add_box(slide, Inches(8.9), y, box_w, box_h, ACCENT_PURPLE,
            "EXTRACTION PATTERNS", 16, WHITE, True)
    add_text(slide, Inches(9.1), y + Inches(1.0), box_w - Inches(0.4), Inches(0.8),
             "Extract field values from text\n• LINE templates (A, 9, literal)\n• FIELD positions (col range)\n• TABLE_DEF for repeating rows", 10, WHITE)

    # Interaction description
    y2 = Inches(3.6)
    add_text(slide, Inches(0.5), y2, Inches(12), Inches(0.5),
             "How they interact at runtime:", 16, DARK_BLUE, True)

    add_text(slide, Inches(0.5), y2 + Inches(0.5), Inches(12), Inches(3.2),
             "1. MAP SEARCH finds all pages matching a field value (report-wide)\n"
             "     → e.g., ACCOUNT_NO = '200-044295-001' → pages [117, 120, 3200]\n\n"
             "2. SECTION SECURITY filters results to user's authorized page ranges\n"
             "     → User authorized for Section '501' (pages 890-3093) → only page 3200\n\n"
             "3. EXTRACTION PATTERNS read the matched pages and extract structured data\n"
             "     → LINE template matching → FIELD extraction at column positions → CSV/JSON output",
             12, DARK_GRAY)

    # Key insight box
    add_box(slide, Inches(0.5), Inches(6.3), Inches(12), Inches(0.9), RGBColor(0xFD, 0xF2, 0xE9),
            "KEY INSIGHT: Indexed fields and sections are independent.\n"
            "MAP search is always report-wide; section permissions are applied as a filter afterwards.",
            12, DARK_GRAY, False, PP_ALIGN.LEFT)


def add_section_security_slide(prs):
    """Slide 3: Section security flow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "1. Section Security — Access Control Flow", 28, DARK_BLUE, True)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(12), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = MED_BLUE
    bar.line.fill.background()

    # Flow diagram
    y = Inches(1.3)

    # User box
    add_box(slide, Inches(0.5), y, Inches(2.2), Inches(0.8), DARK_BLUE,
            "USER\nWindows SID", 12, WHITE, True)

    # Arrow text
    add_text(slide, Inches(2.8), y + Inches(0.15), Inches(1), Inches(0.5), "→", 24, MED_GRAY, True)

    # STYPE_SECTION
    add_box(slide, Inches(3.5), y, Inches(3.2), Inches(0.8), MED_BLUE,
            "STYPE_SECTION\nACL check per SECTION_ID", 11, WHITE, True)

    add_text(slide, Inches(6.8), y + Inches(0.15), Inches(1), Inches(0.5), "→", 24, MED_GRAY, True)

    # Allowed sections
    add_box(slide, Inches(7.5), y, Inches(2.5), Inches(0.8), ACCENT_GREEN,
            "ALLOWED\nSECTION_IDs", 12, WHITE, True)

    add_text(slide, Inches(10.1), y + Inches(0.15), Inches(1), Inches(0.5), "→", 24, MED_GRAY, True)

    # RPT SECTIONHDR
    add_box(slide, Inches(10.5), y, Inches(2.3), Inches(0.8), ACCENT_ORANGE,
            "RPT SECTIONHDR\npage ranges", 11, WHITE, True)

    # Second row: Example
    y2 = Inches(2.8)
    add_text(slide, Inches(0.5), y2, Inches(12), Inches(0.5),
             "Example — DDU017P Report (3500 pages, 3 sections):", 14, DARK_BLUE, True)

    # Section table
    y3 = Inches(3.4)
    headers = ["SECTION_ID", "NAME", "START_PAGE", "PAGE_COUNT", "Page Range"]
    data = [
        ["14259", "501", "1", "890", "1 – 890"],
        ["14260", "201", "891", "2203", "891 – 3093"],
        ["14261", "305", "3094", "407", "3094 – 3500"],
    ]

    col_w = [Inches(1.6), Inches(1.2), Inches(1.6), Inches(1.6), Inches(1.8)]
    x_start = Inches(1.5)

    # Header row
    x = x_start
    for i, h in enumerate(headers):
        add_box(slide, x, y3, col_w[i], Inches(0.4), DARK_BLUE, h, 10, WHITE, True)
        x += col_w[i]

    # Data rows
    for row_idx, row in enumerate(data):
        x = x_start
        y_row = y3 + Inches(0.4) + Inches(0.35) * row_idx
        bg = LIGHT_GRAY if row_idx % 2 == 0 else WHITE
        for i, val in enumerate(row):
            add_box(slide, x, y_row, col_w[i], Inches(0.35), bg, val, 10, BLACK)
            x += col_w[i]

    # User scenarios
    y4 = Inches(5.0)
    add_text(slide, Inches(0.5), y4, Inches(12), Inches(0.4),
             "User Access Scenarios:", 14, DARK_BLUE, True)

    add_text(slide, Inches(0.5), y4 + Inches(0.5), Inches(5.8), Inches(1.5),
             "User A — authorized for section '501' only:\n"
             "  → Sees pages 1-890\n"
             "  → Cannot see pages 891+",
             11, DARK_GRAY)

    add_text(slide, Inches(6.5), y4 + Inches(0.5), Inches(5.8), Inches(1.5),
             "User B — authorized for sections '501' + '305':\n"
             "  → Sees pages 1-890 AND 3094-3500\n"
             "  → Cannot see pages 891-3093 (section '201')",
             11, DARK_GRAY)


def add_map_indexing_slide(prs):
    """Slide 4: MAP file indexing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "2. MAP File Indexing — Field Search Flow", 28, DARK_BLUE, True)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(12), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_GREEN
    bar.line.fill.background()

    # MAP file structure
    y = Inches(1.3)
    add_text(slide, Inches(0.5), y, Inches(5), Inches(0.4),
             "MAP File Structure:", 14, DARK_BLUE, True)

    # MAP segments visual
    y1 = Inches(1.8)
    add_box(slide, Inches(0.5), y1, Inches(5.5), Inches(0.5), DARK_BLUE,
            "MAPHDR — header (signature, segment count, date)", 10, WHITE)

    add_box(slide, Inches(0.5), y1 + Inches(0.6), Inches(5.5), Inches(0.7), MED_BLUE,
            "Segment 0 — Master Lookup\n(LINE_ID, FIELD_ID) → segment#  |  record_id → page#", 10, WHITE)

    add_box(slide, Inches(0.5), y1 + Inches(1.4), Inches(5.5), Inches(0.5), ACCENT_GREEN,
            "Segment 1 — ACCOUNT_NO index (sorted values → u32_index)", 10, WHITE)

    add_box(slide, Inches(0.5), y1 + Inches(2.0), Inches(5.5), Inches(0.5), RGBColor(0x24, 0x72, 0x3E),
            "Segment 2 — CIF_NUMBER index (sorted values → u32_index)", 10, WHITE)

    add_box(slide, Inches(0.5), y1 + Inches(2.6), Inches(5.5), Inches(0.5), RGBColor(0x1A, 0x55, 0x2D),
            "Segment N — ... more IS_INDEXED fields ...", 10, WHITE)

    # Search flow (right side)
    add_text(slide, Inches(6.8), y, Inches(5), Inches(0.4),
             "Search Flow:", 14, DARK_BLUE, True)

    steps = [
        ("1", "User searches:\nACCOUNT_NO = '200-044295-001'", MED_BLUE),
        ("2", "Find ACCOUNT_NO segment\nin MAP file (via Segment 0)", ACCENT_GREEN),
        ("3", "Binary search sorted values\n→ get u32_index = 136,653", RGBColor(0x24, 0x72, 0x3E)),
        ("4", "Join u32_index to Segment 0\n→ page 2748, row 11", MED_BLUE),
        ("5", "Decompress RPT page 2748\nExtract value at row 11", ACCENT_PURPLE),
    ]

    for i, (num, text, color) in enumerate(steps):
        step_y = Inches(1.8) + Inches(0.85) * i
        # Number circle
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.8), step_y, Inches(0.4), Inches(0.4))
        circ.fill.solid()
        circ.fill.fore_color.rgb = color
        circ.line.fill.background()
        tf = circ.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        add_text(slide, Inches(7.3), step_y, Inches(5), Inches(0.7), text, 11, DARK_GRAY)

    # Section intersection note
    add_box(slide, Inches(0.5), Inches(6.2), Inches(12), Inches(1.0), RGBColor(0xE8, 0xF5, 0xE9),
            "SECTION INTERSECTION: MAP search returns pages [117, 120, 3200] (report-wide)\n"
            "User authorized for section '501' (pages 890-3093) → only page 3200 is returned\n"
            "Indexed fields and sections are INDEPENDENT — intersection happens at serve time",
            11, DARK_GRAY, False, PP_ALIGN.LEFT)


def add_extraction_slide(prs):
    """Slide 5: Extraction patterns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "3. Extraction Patterns — LINE Templates & FIELD Positions", 28, DARK_BLUE, True)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(12), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_PURPLE
    bar.line.fill.background()

    # Template matching concept
    y = Inches(1.3)
    add_text(slide, Inches(0.5), y, Inches(12), Inches(0.4),
             "LINE Template Matching — How text lines are classified:", 14, DARK_BLUE, True)

    # Template example
    y1 = Inches(1.8)
    add_box(slide, Inches(0.5), y1, Inches(12), Inches(0.5), RGBColor(0xF5, 0xF5, 0xF5),
            "Report text line:   200-044295-001  JOHN DOE            15/04/2025  1,234.56", 11, BLACK, False, PP_ALIGN.LEFT)

    add_box(slide, Inches(0.5), y1 + Inches(0.55), Inches(12), Inches(0.5), RGBColor(0xE8, 0xE0, 0xF0),
            "LINE.TEMPLATE:      999-999999-999  AAAA AAAA            99/99/9999  9,999.99", 11, ACCENT_PURPLE, False, PP_ALIGN.LEFT)

    # Legend
    add_text(slide, Inches(0.5), y1 + Inches(1.2), Inches(12), Inches(0.6),
             "Template characters:  A = alpha (weight 1.0)  |  9 = digit (weight 1.0)  |  "
             "Literals (-, /, .) = exact match (weight 3.0)  |  space = expected space (weight 0.5)\n"
             "Matching threshold: 0.55 — literal characters act as structural anchors for reliable classification",
             10, MED_GRAY)

    # FIELD extraction
    y2 = Inches(3.5)
    add_text(slide, Inches(0.5), y2, Inches(12), Inches(0.4),
             "FIELD Extraction — Column positions within matched LINE:", 14, DARK_BLUE, True)

    # Visual field positions
    y3 = Inches(4.0)
    fields = [
        ("ACCOUNT_NO", Inches(0.5), Inches(2.5), ACCENT_GREEN, "col 1-15"),
        ("CUSTOMER_NAME", Inches(3.2), Inches(3.0), MED_BLUE, "col 17-36"),
        ("DATE", Inches(6.4), Inches(1.8), ACCENT_ORANGE, "col 38-47"),
        ("AMOUNT", Inches(8.4), Inches(1.5), ACCENT_RED, "col 49-56"),
    ]

    # Text line background
    add_box(slide, Inches(0.5), y3, Inches(10.0), Inches(0.45), RGBColor(0xF5, 0xF5, 0xF5),
            "200-044295-001  JOHN DOE            15/04/2025  1,234.56", 11, BLACK, False, PP_ALIGN.LEFT)

    # Field markers below
    for name, left, width, color, col_text in fields:
        add_box(slide, left, y3 + Inches(0.55), width, Inches(0.35), color,
                f"{name}", 9, WHITE, True)
        add_text(slide, left, y3 + Inches(0.92), width, Inches(0.3),
                 col_text, 8, MED_GRAY, False, PP_ALIGN.CENTER)

    # FIELD flags
    y4 = Inches(5.3)
    add_text(slide, Inches(0.5), y4, Inches(12), Inches(0.4),
             "FIELD Flags — Same positions serve multiple purposes:", 14, DARK_BLUE, True)

    flags = [
        ("IS_INDEXED = 1", "Value stored in MAP file for search\n(e.g., ACCOUNT_NO → MAP segment)", ACCENT_GREEN),
        ("IS_SIGNIFICANT = 1", "Value defines section boundary\n(e.g., BRANCH_CODE → SECTION.NAME)", MED_BLUE),
        ("START/END_COLUMN", "Position-based value extraction\n(runtime field extraction)", ACCENT_PURPLE),
    ]

    for i, (flag, desc, color) in enumerate(flags):
        x = Inches(0.5) + Inches(4.1) * i
        add_box(slide, x, y4 + Inches(0.5), Inches(3.8), Inches(0.4), color, flag, 11, WHITE, True)
        add_text(slide, x + Inches(0.1), y4 + Inches(1.0), Inches(3.6), Inches(0.6), desc, 10, DARK_GRAY)


def add_table_extraction_slide(prs):
    """Slide 6: TABLE_DEF/TABLE_ITEM extraction."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "4. Table Extraction — Multi-Row Structures", 28, DARK_BLUE, True)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(12), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_ORANGE
    bar.line.fill.background()

    # TABLE_DEF concept
    y = Inches(1.3)
    add_text(slide, Inches(0.5), y, Inches(12), Inches(0.4),
             "TABLE_DEF / TABLE_ITEM — Grouping LINEs into repeating structures:", 14, DARK_BLUE, True)

    # Visual: table structure
    y1 = Inches(1.9)

    # Left: TABLE_DEF
    add_box(slide, Inches(0.5), y1, Inches(3.5), Inches(0.5), DARK_BLUE,
            "TABLE_DEF: 'Transaction Details'", 11, WHITE, True)

    # TABLE_ITEM entries
    items = [
        ("LINE_ID 10 — Header", "ACCOUNT  NAME  DATE  AMOUNT", MED_BLUE),
        ("LINE_ID 11 — Detail (repeating)", "200-044295-001  JOHN DOE  15/04  1,234.56", ACCENT_GREEN),
        ("LINE_ID 11 — Detail (repeating)", "200-044295-002  JANE DOE  16/04  2,567.89", ACCENT_GREEN),
        ("LINE_ID 11 — Detail (repeating)", "300-012345-001  BOB SMITH  17/04  890.12", ACCENT_GREEN),
        ("LINE_ID 12 — Subtotal", "                    TOTAL:  4,692.57", ACCENT_ORANGE),
    ]

    for i, (line_type, example, color) in enumerate(items):
        yi = y1 + Inches(0.6) + Inches(0.55) * i
        add_box(slide, Inches(0.5), yi, Inches(3.5), Inches(0.5), color,
                line_type, 9, WHITE, True)
        add_box(slide, Inches(4.2), yi, Inches(7.5), Inches(0.5), RGBColor(0xF8, 0xF8, 0xF8),
                example, 10, BLACK, False, PP_ALIGN.LEFT)

    # Right: Extraction pipeline
    y2 = Inches(4.8)
    add_text(slide, Inches(0.5), y2, Inches(12), Inches(0.4),
             "Complete Extraction Pipeline:", 14, DARK_BLUE, True)

    steps = [
        ("1. Resolve", "Report name → REPORT_SPECIES_ID → REPORT_INSTANCE\n→ STRUCTURE_DEF_ID, MAP file, RPT file", MED_BLUE),
        ("2. Search", "ACCOUNT_NO value → MAP binary search → page numbers", ACCENT_GREEN),
        ("3. Decompress", "Page number → PAGETBLHDR offset → zlib → raw text lines", ACCENT_ORANGE),
        ("4. Classify", "Each text line → score against all LINE templates → best match", ACCENT_PURPLE),
        ("5. Extract", "Matched line → FIELD at [START_COL:END_COL] → field values", ACCENT_RED),
        ("6. Output", "Structured CSV/JSON with all extracted fields", DARK_BLUE),
    ]

    x_pos = Inches(0.5)
    for i, (step, desc, color) in enumerate(steps):
        xi = x_pos + Inches(2.1) * i
        add_box(slide, xi, y2 + Inches(0.5), Inches(1.9), Inches(0.5), color, step, 10, WHITE, True)
        add_text(slide, xi, y2 + Inches(1.1), Inches(2.0), Inches(1.0), desc, 8, DARK_GRAY)


def add_complete_flow_slide(prs):
    """Slide 7: Complete data flow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "5. Complete Data Flow — Ingestion to Access", 28, DARK_BLUE, True)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(12), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    # Ingestion column
    y = Inches(1.3)
    add_text(slide, Inches(0.5), y, Inches(4), Inches(0.4), "INGESTION", 16, DARK_BLUE, True)

    add_box(slide, Inches(0.5), y + Inches(0.5), Inches(3.8), Inches(0.5), MED_BLUE,
            "Spool arrives → SIGNATURE match", 10, WHITE)
    add_box(slide, Inches(0.5), y + Inches(1.1), Inches(3.8), Inches(0.5), MED_BLUE,
            "LINE template → classify each line", 10, WHITE)
    add_box(slide, Inches(0.5), y + Inches(1.7), Inches(3.8), Inches(0.5), ACCENT_GREEN,
            "IS_INDEXED fields → MAP file segments", 10, WHITE)
    add_box(slide, Inches(0.5), y + Inches(2.3), Inches(3.8), Inches(0.5), ACCENT_ORANGE,
            "IS_SIGNIFICANT fields → SECTION boundaries", 10, WHITE)
    add_box(slide, Inches(0.5), y + Inches(2.9), Inches(3.8), Inches(0.5), ACCENT_PURPLE,
            "SENSITIVE_FIELD → report name, date, sections", 10, WHITE)

    # Storage column
    add_text(slide, Inches(4.8), y, Inches(3.5), Inches(0.4), "STORAGE", 16, DARK_BLUE, True)

    add_box(slide, Inches(4.8), y + Inches(0.5), Inches(3.5), Inches(0.5), DARK_BLUE,
            "REPORT_INSTANCE record", 10, WHITE)
    add_box(slide, Inches(4.8), y + Inches(1.3), Inches(3.5), Inches(1.0), ACCENT_ORANGE,
            "RPT file\n• Compressed pages (zlib)\n• SECTIONHDR (section→pages)\n• PAGETBLHDR (page offsets)", 9, WHITE)
    add_box(slide, Inches(4.8), y + Inches(2.5), Inches(3.5), Inches(0.9), ACCENT_GREEN,
            "MAP file\n• Segment 0: master index\n• Segments 1-N: sorted field indices", 9, WHITE)

    # Access column
    add_text(slide, Inches(8.8), y, Inches(4), Inches(0.4), "USER ACCESS", 16, DARK_BLUE, True)

    add_box(slide, Inches(8.8), y + Inches(0.5), Inches(3.8), Inches(0.5), MED_BLUE,
            "1. User auth → Windows SID", 10, WHITE)
    add_box(slide, Inches(8.8), y + Inches(1.1), Inches(3.8), Inches(0.5), MED_BLUE,
            "2. STYPE_SECTION → allowed sections", 10, WHITE)
    add_box(slide, Inches(8.8), y + Inches(1.7), Inches(3.8), Inches(0.5), ACCENT_GREEN,
            "3. MAP search → matching pages", 10, WHITE)
    add_box(slide, Inches(8.8), y + Inches(2.3), Inches(3.8), Inches(0.5), ACCENT_ORANGE,
            "4. Intersect pages ∩ allowed sections", 10, WHITE)
    add_box(slide, Inches(8.8), y + Inches(2.9), Inches(3.8), Inches(0.5), ACCENT_PURPLE,
            "5. Decompress + extract → serve to user", 10, WHITE)

    # Existing code status
    y3 = Inches(4.8)
    add_text(slide, Inches(0.5), y3, Inches(12), Inches(0.4),
             "Implementation Status:", 16, DARK_BLUE, True)

    done = [
        ("intellistor_extractor.py", "LINE matching, FIELD extraction, MAP search, RPT decompress"),
        ("intellistor_viewer.py", "Database access, MAP file parser, data classes"),
        ("rpt_page_extractor.py", "Page decompression, PAGETBLHDR, binary objects"),
        ("extract_instances_sections.py", "Instance CSV extraction with section/MAP support"),
    ]
    todo = [
        ("TABLE_DEF/TABLE_ITEM", "Multi-row table extraction (repeating detail rows)"),
        ("FOLLOW_LIST", "Line sequence validation (disambiguation)"),
        ("FIELD_TYPE", "Formatting/parsing rules (date, numeric, locale)"),
    ]

    for i, (name, desc) in enumerate(done):
        yi = y3 + Inches(0.5) + Inches(0.35) * i
        add_box(slide, Inches(0.5), yi, Inches(0.3), Inches(0.3), ACCENT_GREEN, "✓", 12, WHITE, True)
        add_text(slide, Inches(0.9), yi, Inches(4), Inches(0.3), name, 10, DARK_GRAY, True)
        add_text(slide, Inches(4.5), yi, Inches(7), Inches(0.3), desc, 10, MED_GRAY)

    for i, (name, desc) in enumerate(todo):
        yi = y3 + Inches(0.5) + Inches(0.35) * (len(done) + i)
        add_box(slide, Inches(0.5), yi, Inches(0.3), Inches(0.3), ACCENT_RED, "○", 12, WHITE, True)
        add_text(slide, Inches(0.9), yi, Inches(4), Inches(0.3), name, 10, DARK_GRAY, True)
        add_text(slide, Inches(4.5), yi, Inches(7), Inches(0.3), desc, 10, MED_GRAY)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_overview_slide(prs)
    add_section_security_slide(prs)
    add_map_indexing_slide(prs)
    add_extraction_slide(prs)
    add_table_extraction_slide(prs)
    add_complete_flow_slide(prs)

    output = '/Volumes/acasis/projects/python/ocbc/IntelliSTOR_Migration/4_Migration_Instances/IntelliSTOR_Architecture.pptx'
    prs.save(output)
    print(f'PowerPoint saved to: {output}')


if __name__ == '__main__':
    main()
